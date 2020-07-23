import tempfile
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor


def create(filename, prs_content, background_color, text_color):
    if (not ('.pptx' in filename)):
        filename += '.pptx'

    prs = Presentation()

    content_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(content_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(
        background_color[0], background_color[1], background_color[2])

    left = top = Inches(1)
    width = Inches(8)
    height = Inches(6)

    content = slide.shapes.add_textbox(left, top, width, height)
    content.text = prs_content

    # most of the shapes only contain one paragraph for this use case
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(
        text_color[0], text_color[1], text_color[2])

    pptx_path = tempfile.gettempdir() + '/' + filename
    prs.save(pptx_path)
    return pptx_path
