import tempfile
from pptx import Presentation
from pptx.util import Inches


def create(filename, prs_content):
    prs = Presentation()

    content_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(content_slide_layout)

    left = top = Inches(1)
    width = Inches(8)
    height = Inches(6)

    content = slide.shapes.add_textbox(left, top, width, height)

    content.text = prs_content
    pptx_path = tempfile.gettempdir() + '/' + filename
    prs.save(pptx_path)
    return pptx_path
