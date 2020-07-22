from pptx import Presentation
from pptx.util import Inches


def create_pptx(prs_content, file_name):
    prs = Presentation()

    content_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(content_slide_layout)

    left = top = Inches(1)
    width = Inches(8)
    height = Inches(6)

    content = slide.shapes.add_textbox(left, top, width, height)

    content.text = prs_content
    prs.save(file_name)


